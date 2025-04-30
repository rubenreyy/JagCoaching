import os
import json
import io
import base64
import re
import logging
from pptx import Presentation
from dotenv import load_dotenv
import google.generativeai as genai
from PIL import Image

# Configure logging
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

# Load environment variables
load_dotenv(dotenv_path=".env.development")

# Configure Gemini
gemini_api_key = os.getenv("GOOGLE_GEMINI_API_KEY")
genai.configure(api_key=gemini_api_key)
model = genai.GenerativeModel("gemini-1.5-pro-latest") 

def extract_presentation_content(file_path):
    """Extract text and image content from a PowerPoint presentation."""
    logger.info(f"Extracting content from presentation: {file_path}")
    presentation = Presentation(file_path)
    slides_content = []

    for slide_number, slide in enumerate(presentation.slides, start=1):
        slide_data = {"slide_number": slide_number, "text": [], "images": []}

        for shape in slide.shapes:
            # Extract text
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    slide_data["text"].append(paragraph.text)

            # Extract images
            if shape.shape_type == 13:  # 13 = Picture
                image = shape.image
                image_bytes = image.blob
                image_stream = io.BytesIO(image_bytes)
                try:
                    pil_image = Image.open(image_stream).convert("RGB")
                    slide_data["images"].append(pil_image)
                except Exception as e:
                    logger.warning(f"Unable to decode image on slide {slide_number}: {e}")

        slides_content.append(slide_data)

    logger.info(f"Extracted content from {len(slides_content)} slides.")
    return slides_content

def analyze_presentation_with_gemini(slides_content):
    #Send slide content and images to Gemini for analysis.
    logger.info("Preparing multimodal prompt for Gemini")

    prompt_text = """
You are an expert presentation coach. Analyze this PowerPoint presentation and provide feedback on:
1. Slide content clarity
2. Visual design quality (based on image layout and text/image use)
3. Key topics and themes
4. Suggestions for improvement

Respond in the following JSON format:
{
  "slide_clarity": "Very brief feedback about the clarity of the slide content",
  "visual_quality": "Very brief feedback about the quality and choices of visual design/how well they relate to the topic",
  "key_topics": "Very brief feedback about they key topics and themes of the presentation",
  "overall_suggestion": "One concise improvement tip"
}
"""

    # Gemini accepting as input chunks
    input_parts = [prompt_text]

    for slide in slides_content:
        slide_text = f"\nSlide {slide['slide_number']} content:\n" + "\n".join(slide["text"])
        input_parts.append(slide_text)
        for image in slide["images"]:
            input_parts.append(image)

    try:
        logger.info("Sending multimodal request to Gemini...")
        response = model.generate_content(input_parts)
        response_text = response.text

        # Extract JSON from markdown formatting if needed
        json_match = re.search(r'```json\s*(\{.*?\})\s*```', response_text, re.DOTALL)
        cleaned_json = json_match.group(1) if json_match else response_text

        try:
            result = json.loads(cleaned_json)
            logger.info("Successfully parsed JSON response.")
            return result
        except json.JSONDecodeError as e:
            logger.error(f"Failed to parse JSON: {e}")
            return {"raw_response": response_text}

    except Exception as e:
        logger.exception("Gemini Vision analysis failed.")
        return {"error": str(e)}

def analyze_presentation(file_path):
    """Main function to analyze a PowerPoint presentation."""
    slides_content = extract_presentation_content(file_path)
    analysis_result = analyze_presentation_with_gemini(slides_content)
    return analysis_result

# Example usage
if __name__ == "__main__":
    '''presentation_file = "path/to/your/ppt.pptx"
    if os.path.exists(presentation_file):
        logger.info(f"Analyzing presentation: {presentation_file}")
        result = analyze_presentation(presentation_file)
        print(json.dumps(result, indent=4))
    else:
        logger.error(f"Presentation file not found: {presentation_file}")'''
